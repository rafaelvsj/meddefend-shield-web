from pptx import Presentation
import logging
from typing import Dict, Any
from .base_extractor import BaseExtractor

logger = logging.getLogger(__name__)

class PPTXExtractor(BaseExtractor):
    """Extract text from PPTX files using python-pptx."""
    
    async def extract(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Extract text from PPTX slides in order."""
        
        try:
            prs = Presentation(file_path)
            slides_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num} ---")
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_data.append(cell.text.strip())
                            if row_data:
                                slide_text.append(" | ".join(row_data))
                
                if len(slide_text) > 1:  # More than just the slide header
                    slides_content.append("\n".join(slide_text))
            
            full_text = "\n\n".join(slides_content)
            
            if len(full_text.strip()) < 10:
                raise Exception("No text content found in PPTX")
            
            logger.info(f"Extracted {len(full_text)} characters from {len(slides_content)} slides")
            
            return {
                'text': full_text,
                'method': 'python-pptx',
                'ocr_used': False
            }
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
            raise Exception(f"PPTX extraction failed: {str(e)}")